import os
import json
import pdfplumber
import pandas as pd
import requests
from pptx import Presentation
from docx import Document
from urllib.parse import urlparse, urljoin
import xml.etree.ElementTree as ET
from sentence_transformers import SentenceTransformer
import psycopg2
from psycopg2.extras import execute_values
from pgvector.psycopg2 import register_vector
from concurrent.futures import ThreadPoolExecutor, as_completed

# === Configuration ===
DB_CONFIG = {
    "dbname": "market_db",
    "user": "postgres",
    "password": "1329",
    "host": "localhost",
    "port": "5432"
}
TABLE_NAME = "your_table"  # Ensure table exists with (id SERIAL, content TEXT, embedding VECTOR(384))

# === Connect to DB ===
conn = psycopg2.connect(**DB_CONFIG)
register_vector(conn)

# === Load Embedder ===
embedder = SentenceTransformer("all-MiniLM-L6-v2")

# === Utility: Split text ===
def split_text(text, chunk_size=300):
    words = text.split()
    return [" ".join(words[i:i + chunk_size]) for i in range(0, len(words), chunk_size)]

# === Extractors ===
def extract_text_from_pdf(file_path):
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs)

def extract_text_from_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_csv(file_path):
    df = pd.read_csv(file_path)
    return "\n".join(df.astype(str).apply(lambda row: " ".join(row), axis=1))

def extract_text_from_json(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    return json.dumps(data, indent=2)

def extract_text_from_url(url):
    response = requests.get(url, timeout=10)
    response.raise_for_status()
    return response.text

# === Sitemap URL Extractor ===
def get_urls_from_sitemap(domain):
    sitemap_url = urljoin(domain, "/sitemap.xml")
    try:
        response = requests.get(sitemap_url, timeout=10)
        response.raise_for_status()
        root = ET.fromstring(response.content)
        urls = [elem.text for elem in root.iter() if elem.tag.endswith("loc")]
        print(f"🔍 Found {len(urls)} URLs in sitemap: {sitemap_url}")
        return urls
    except Exception as e:
        print(f"⚠️ Failed to fetch sitemap from {sitemap_url}: {e}")
        return []

# === Main ingestion function with batch DB insert ===
def ingest_file(file_path_or_url):
    ext = ""
    text = ""

    if file_path_or_url.startswith("http"):
        text = extract_text_from_url(file_path_or_url)
    else:
        ext = os.path.splitext(file_path_or_url)[1].lower()

        if ext == ".pdf":
            text = extract_text_from_pdf(file_path_or_url)
        elif ext == ".docx":
            text = extract_text_from_docx(file_path_or_url)
        elif ext == ".txt":
            text = extract_text_from_txt(file_path_or_url)
        elif ext == ".pptx":
            text = extract_text_from_pptx(file_path_or_url)
        elif ext == ".csv":
            text = extract_text_from_csv(file_path_or_url)
        elif ext == ".json":
            text = extract_text_from_json(file_path_or_url)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    # Chunk text
    chunks = split_text(text)
    # Generate embeddings in batch
    embeddings = embedder.encode(chunks).tolist()

    records = [(chunk, embedding) for chunk, embedding in zip(chunks, embeddings)]

    with conn.cursor() as cur:
        execute_values(
            cur,
            f"INSERT INTO {TABLE_NAME} (content, embedding) VALUES %s",
            records
        )
        conn.commit()

    print(f"✅ Ingested {len(chunks)} chunks from: {file_path_or_url}")

# === Auto-detect Domain / URL / File and ingest ===
def auto_ingest(input_str):
    if input_str.startswith("http") and (input_str.endswith(".xml") or input_str.endswith(".sitemap")):
        # Direct sitemap URL
        urls = get_urls_from_sitemap(input_str)
        parallel_ingest(urls)

    elif input_str.startswith("http") and "." not in os.path.splitext(urlparse(input_str).path)[1]:
        # Domain, try sitemap.xml
        urls = get_urls_from_sitemap(input_str)
        parallel_ingest(urls)

    else:
        # Single file or page
        try:
            ingest_file(input_str)
        except Exception as e:
            print(f"❌ Failed to ingest {input_str}: {e}")

# === Parallel ingestion of multiple URLs ===
def parallel_ingest(urls, max_workers=5):
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = [executor.submit(ingest_file, url) for url in urls]
        for future in as_completed(futures):
            try:
                future.result()
            except Exception as e:
                print(f"❌ Error during parallel ingestion: {e}")

# === Example Usage ===
if __name__ == "__main__":
    test_inputs = [
        "https://introspectivemarketresearch.com/sitemap.xml"  # Just domain - will fetch sitemap and ingest URLs in parallel
        # "https://www.example.com/sitemap.xml"     # Direct sitemap URL example
        # "https://www.databridgemarketresearch.com/reports/global-rice-husk-ash-market"  # Single page URL example
        # "my_local_file.pdf"                        # Local file example
    ]

    for item in test_inputs:
        auto_ingest(item)
